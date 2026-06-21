"""
Generate a 10-12 slide presentation for reporting on the
LLM system infrastructure and acceleration techniques.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

ROOT = Path(__file__).resolve().parents[1]
CHARTS = ROOT / "docs" / "charts"
OUT = ROOT / "docs" / "LLM_Infra_Acceleration_Report.pptx"

# ── Color Palette ────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1a, 0x1a, 0x2e)
BG_MID    = RGBColor(0x22, 0x22, 0x3a)
ACCENT    = RGBColor(0x00, 0xd2, 0xff)  # cyan
ACCENT2   = RGBColor(0xff, 0x6b, 0x6b)  # coral
WHITE     = RGBColor(0xff, 0xff, 0xff)
GRAY      = RGBColor(0xaa, 0xaa, 0xbb)
GREEN     = RGBColor(0x00, 0xc9, 0x7b)
ORANGE    = RGBColor(0xff, 0xa5, 0x00)
LIGHT_BG  = RGBColor(0xf8, 0xf9, 0xfa)

# Use light theme for readability
SLIDE_BG  = RGBColor(0xff, 0xff, 0xff)
TITLE_CLR = RGBColor(0x1a, 0x1a, 0x2e)
TEXT_CLR  = RGBColor(0x33, 0x33, 0x33)
ACCENT_BAR= RGBColor(0x00, 0x6d, 0xfd)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def _set_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_accent_bar(slide, top=Inches(1.3), width=Inches(1.5)):
    """Add a colored accent bar under the title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, width, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BAR
    shape.line.fill.background()


def _title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, RGBColor(0x0d, 0x1b, 0x2a))

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(0x88, 0xcc, 0xff)

    # Bottom line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.5), Inches(3), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    return slide


def _content_slide(title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_CLR

    _add_accent_bar(slide)

    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        # Support bold prefix with "**text**" syntax
        if "**" in bullet:
            parts = bullet.split("**")
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                run.text = part
                run.font.size = Pt(20)
                run.font.color.rgb = TEXT_CLR
                if j % 2 == 1:  # odd parts are bold
                    run.font.bold = True
        else:
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_CLR

        p.space_after = Pt(12)

    if note:
        txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.7))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(14)
        p3.font.color.rgb = GRAY
        p3.font.italic = True

    return slide


def _chart_slide(title, chart_filename, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE_CLR

    _add_accent_bar(slide, top=Inches(1.1))

    # Chart image
    chart_path = CHARTS / chart_filename
    if chart_path.exists():
        slide.shapes.add_picture(
            str(chart_path), Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.2)
        )
    else:
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"[Chart: {chart_filename} not found]"
        p2.font.size = Pt(18)
        p2.font.color.rgb = ACCENT2

    if caption:
        txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11), Inches(0.6))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = caption
        p3.font.size = Pt(14)
        p3.font.color.rgb = GRAY
        p3.font.italic = True

    return slide


def _table_slide(title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE_CLR

    _add_accent_bar(slide, top=Inches(1.1))

    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        total = 11.5
        col_widths = [total / n_cols] * n_cols

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.8), Inches(1.5),
        Inches(sum(col_widths)), Inches(0.5 * n_rows)
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BAR
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(15)
                para.font.color.rgb = TEXT_CLR
                para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Alternate row colors
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xf0, 0xf4, 0xf8) if i % 2 == 0 else WHITE

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

# ── Slide 1: Title ───────────────────────────────────────────────────────────
_title_slide(
    "LLM Infrastructure Acceleration\nfor Financial QA",
    "Local-First RAG System  |  Llama-3.2-3B  |  vLLM + AWQ4 Quantization\nRTX 3090 Ti  |  May 2026"
)

# ── Slide 2: Problem & Motivation ────────────────────────────────────────────
_content_slide(
    "Problem & Motivation",
    [
        "Financial analysts need fast, accurate answers from SEC filings (10-K, 10-Q, 8-K)",
        "Cloud LLM APIs (GPT-4o) have **knowledge cutoff issues** -- miss recent FY2023 data",
        "Cloud APIs **hallucinate financial numbers** -- 42.9% value accuracy on our benchmark",
        "Privacy concerns: sending proprietary financial queries to external APIs",
        "Goal: **local-first system** on a single GPU that beats cloud accuracy at lower cost",
    ],
    note="92 S&P 100 tickers | 2019-2024 | ~120K text chunks + ~40K structured XBRL rows"
)

# ── Slide 3: System Architecture ─────────────────────────────────────────────
_content_slide(
    "System Architecture",
    [
        "**3-stage pipeline:** Intent Classification -> Retrieval -> Answer Generation",
        "**Intent Router** (LoRA adapter): classifies query as Type1/Type2/Type3",
        "  Type1 (exact fact) -> NL2SQL -> SQLite (XBRL financial data)",
        "  Type2 (qualitative) -> Query Rewrite -> Hybrid RAG (BM25 + Dense + Cross-Encoder)",
        "  Type3 (chat) -> Direct answer, no retrieval",
        "**vLLM + Punica Multi-LoRA:** 3 adapters share one GPU process",
        "**Data stores:** ChromaDB (516K chunks) + SQLite (40K financial rows)",
    ],
    note="All 3 LoRA adapters (intent, rewriter, nl2sql) served via Punica SGMV kernels in a single vLLM process"
)

# ── Slide 4: Quantization Strategy ───────────────────────────────────────────
_table_slide(
    "Quantization Strategy: FP16 vs INT8 vs AWQ4",
    ["Method", "Technique", "Weight Size", "VRAM (weights)", "KV Cache Budget"],
    [
        ["FP16",  "bfloat16, no compression",           "16-bit", "~6 GB",   "~18 GB"],
        ["INT8",  "bitsandbytes LLM.int8() dequant",    "8-bit",  "~3 GB",   "~21 GB"],
        ["AWQ4",  "W4A16 Marlin fused INT4xFP16 GEMM",  "4-bit",  "~1.5 GB", "~22.5 GB"],
    ],
    col_widths=[1.5, 4.0, 1.5, 2.0, 2.5],
)

# ── Slide 5: Accuracy Comparison ─────────────────────────────────────────────
_chart_slide(
    "Accuracy: Unified 52-Case Benchmark",
    "accuracy_comparison.png",
    caption="AWQ4 achieves the highest overall composite score (0.904), crossing the 0.9 target. "
            "GPT-4o fails on Type1 due to knowledge cutoff."
)

# ── Slide 6: Accuracy Detail ─────────────────────────────────────────────────
_chart_slide(
    "Accuracy Detail: Intent / Value / Keyword",
    "accuracy_detail.png",
    caption="All local configs achieve 100% value accuracy (numbers from SQLite/XBRL). "
            "GPT-4o has best fluency (4.7/5) but worst value accuracy (42.9%)."
)

# ── Slide 7: Throughput ──────────────────────────────────────────────────────
_chart_slide(
    "Throughput: QPS vs Concurrency",
    "throughput_comparison.png",
    caption="AWQ4 is the best local option at every concurrency level. "
            "INT8 degrades severely at c=4 (0.19 QPS) due to runtime dequantization overhead."
)

# ── Slide 8: Latency ─────────────────────────────────────────────────────────
_chart_slide(
    "Latency: Median Response Time vs Concurrency",
    "latency_comparison.png",
    caption="AWQ4 maintains flat ~2.6s p50 at c=4-16. INT8 spikes to 24.6s at c=8. "
            "GPT-4o has lowest latency but worst accuracy."
)

# ── Slide 9: Pipeline Latency Breakdown ──────────────────────────────────────
_chart_slide(
    "End-to-End Pipeline Latency",
    "pipeline_latency.png",
    caption="AWQ4 p50 = 0.8s (Type1 queries resolve in <1s via SQL). "
            "INT8 p95 = 71.9s due to bitsandbytes JIT compilation on early requests."
)

# ── Slide 10: Key Optimizations ──────────────────────────────────────────────
_content_slide(
    "Key Optimizations Applied",
    [
        "**AWQ4 W4A16 Marlin kernels** -- 4x weight compression, fused INT4xFP16 GEMM, zero accuracy loss",
        "**Punica Multi-LoRA** -- 3 adapters in 1 GPU process via SGMV batched forward pass",
        "**Cross-encoder on CPU** -- prevents GPU contention with vLLM; ~200ms for 50 candidates",
        "**Rerank candidate cap (top-50)** -- limits cross-encoder work, prevents stalls at high concurrency",
        "**BM25 index built once** at startup (516K docs), reused across all queries",
        "**Domain-specific AWQ calibration** -- 128 financial conversation samples for quantization",
    ],
    note="Speculative decoding with 1B draft model was tested but HURT performance 25x due to HBM bandwidth contention on 3B target"
)

# ── Slide 11: Results Summary ────────────────────────────────────────────────
_table_slide(
    "Final Results: AWQ4 vs Cloud (GPT-4o)",
    ["Metric", "AWQ4 (Local)", "GPT-4o (API)"],
    [
        ["Composite Accuracy",   "0.9038",       "0.6467"],
        ["Value Accuracy",       "100%",          "42.9%"],
        ["Intent Accuracy",      "100%",          "100%"],
        ["Keyword Hit Rate",     "93.9%",         "75.7%"],
        ["Fluency (1-5)",        "3.4",           "4.7"],
        ["Latency p50 (c=1)",   "0.53s",         "0.85s"],
        ["Peak QPS",             "3.46 (c=16)",   "7.03 (c=8)"],
        ["Cost per 1K queries",  "~$0",           "~$15-30"],
        ["Data Privacy",         "Fully local",   "Sent to OpenAI"],
    ],
    col_widths=[4.0, 3.75, 3.75],
)

# ── Slide 12: Where to Improve & Conclusion ──────────────────────────────────
_content_slide(
    "Where to Improve & Next Steps",
    [
        "**Type2 RAG latency (5-11s)** -- main bottleneck; async parallel retrieval + lighter reranker",
        "**Fluency gap** (3.4 vs 4.7) -- fine-tune answer generator on GPT-4o-distilled data",
        "**Streaming** -- add SSE via vLLM async iterator for perceived latency improvement",
        "**Context window** -- upgrade to 32K+ model (Llama-3.1-8B) for more retrieved chunks",
        "**Multi-GPU scaling** -- tensor parallelism or load-balanced vLLM for >5 QPS",
        "**Evaluation coverage** -- expand from 52 to 200+ cases with adversarial examples",
        "",
        "**Conclusion:** A local 3B model with AWQ4 + vLLM achieves 0.90 accuracy,",
        "100% value accuracy, and 0.53s latency -- beating GPT-4o on financial QA at $0/query.",
    ],
)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Presentation saved: {OUT}")
print(f"  {len(prs.slides)} slides")
